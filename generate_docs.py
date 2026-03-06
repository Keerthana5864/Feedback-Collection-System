"""
Faculty Feedback System — Documentation Generator
Generates: PRESENTATION.pptx  and  DOCUMENTATION_PDF.html (open in browser → Print → Save as PDF)
Run:  python generate_docs.py
"""

# ─── PPTX Generation ────────────────────────────────────────────────────────

def generate_pptx():
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        import pptx.util as util
    except ImportError:
        print("Installing python-pptx...")
        import subprocess, sys
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    INDIGO   = RGBColor(0x63, 0x66, 0xF1)
    DARK     = RGBColor(0x1E, 0x29, 0x3B)
    WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
    MUTED    = RGBColor(0x64, 0x74, 0x8B)
    AMBER    = RGBColor(0xFB, 0xBF, 0x24)
    GREEN    = RGBColor(0x10, 0xB9, 0x81)
    RED      = RGBColor(0xEF, 0x44, 0x44)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]   # completely blank layout

    def add_rect(slide, l, t, w, h, fill_rgb):
        shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
        shape.line.fill.background()
        return shape

    def add_text(slide, text, l, t, w, h, size=18, bold=False, color=None, align=PP_ALIGN.LEFT, wrap=True):
        txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        txBox.word_wrap = wrap
        tf = txBox.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color if color else DARK
        return txBox

    def title_slide(title, subtitle, institution):
        slide = prs.slides.add_slide(blank)
        # background
        add_rect(slide, 0, 0, 13.33, 7.5, DARK)
        # indigo accent bar
        add_rect(slide, 0, 5.8, 13.33, 0.12, INDIGO)
        # title
        add_text(slide, title, 1.2, 1.8, 11, 1.4, size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # subtitle
        add_text(slide, subtitle, 1.2, 3.4, 11, 0.7, size=20, color=INDIGO, align=PP_ALIGN.CENTER)
        # institution
        add_text(slide, institution, 1.2, 4.2, 11, 0.5, size=14, color=MUTED, align=PP_ALIGN.CENTER)
        add_text(slide, "February 2026", 1.2, 6.2, 11, 0.5, size=12, color=MUTED, align=PP_ALIGN.CENTER)

    def content_slide(heading, bullets, icon=""):
        slide = prs.slides.add_slide(blank)
        # top bar
        add_rect(slide, 0, 0, 13.33, 1.1, DARK)
        add_rect(slide, 0, 1.1, 0.08, 6.4, INDIGO)
        # heading
        add_text(slide, f"{icon}  {heading}", 0.3, 0.15, 12, 0.8, size=24, bold=True, color=WHITE)
        # bullets
        y = 1.35
        for b in bullets:
            if b.startswith("##"):          # section sub-heading
                add_text(slide, b[2:].strip(), 0.5, y, 12.5, 0.4, size=14, bold=True, color=INDIGO)
                y += 0.45
            elif b.startswith("---"):       # divider — skip
                y += 0.15
            else:
                prefix = "  •  " if not b.startswith("   ") else "      –  "
                add_text(slide, prefix + b.strip(), 0.5, y, 12.5, 0.38, size=13, color=DARK)
                y += 0.4
            if y > 7.1:
                break
        return slide

    def two_col_slide(heading, left_items, right_items, icon=""):
        slide = prs.slides.add_slide(blank)
        add_rect(slide, 0, 0, 13.33, 1.1, DARK)
        add_rect(slide, 0, 1.1, 0.08, 6.4, INDIGO)
        add_text(slide, f"{icon}  {heading}", 0.3, 0.15, 12, 0.8, size=24, bold=True, color=WHITE)
        add_rect(slide, 0.5, 1.25, 5.9, 5.8, RGBColor(0xF8, 0xFA, 0xFC))
        add_rect(slide, 6.9, 1.25, 5.9, 5.8, RGBColor(0xF0, 0xF2, 0xF5))
        y = 1.45
        for item in left_items:
            add_text(slide, "• " + item, 0.7, y, 5.5, 0.38, size=12, color=DARK)
            y += 0.42
        y = 1.45
        for item in right_items:
            add_text(slide, "• " + item, 7.1, y, 5.5, 0.38, size=12, color=DARK)
            y += 0.42

    def thank_you_slide():
        slide = prs.slides.add_slide(blank)
        add_rect(slide, 0, 0, 13.33, 7.5, DARK)
        add_rect(slide, 4, 3.5, 5.33, 0.08, INDIGO)
        add_text(slide, "Thank You", 0, 2.2, 13.33, 1.4, size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, "Faculty Feedback System  |  Your Institution", 0, 3.8, 13.33, 0.5,
                 size=16, color=INDIGO, align=PP_ALIGN.CENTER)
        add_text(slide, "For queries: admin@demo.com", 0, 4.5, 13.33, 0.4,
                 size=13, color=MUTED, align=PP_ALIGN.CENTER)

    # ── Build Slides ──────────────────────────────────────────────────────────

    title_slide(
        "Faculty Feedback System",
        "A digital platform for collecting & analysing faculty feedback",
        "Your Institution  |  Admin Dashboard Project"
    )

    content_slide("Project Overview", [
        "A full-stack web application for your institution",
        "Students submit star-rated feedback for faculty members",
        "Administrators view analytics, alerts & manage responses",
        "Supports secure login, registration & OTP-based password reset",
        "---",
        "## Key Objectives",
        "Enable students to rate faculty on a 1–5 star scale",
        "Give admin real-time analytics via a rich dashboard",
        "Provide instant low-rating alerts to detect issues early",
        "Secure role-based access for students and admins",
    ], icon="📌")

    content_slide("Technology Stack", [
        "## Frontend",
        "HTML5, CSS3, Vanilla JavaScript — no heavy frameworks",
        "Chart.js for bar and pie chart visualisations",
        "Google Fonts (Outfit) for modern typography",
        "---",
        "## Backend",
        "Python 3 + Flask — lightweight REST API server",
        "Flask-CORS for cross-origin request handling",
        "python-dotenv for environment variable management",
        "---",
        "## Database",
        "Supabase (PostgreSQL) — cloud-hosted, real-time capable",
        "Accessed via official Supabase Python client",
    ], icon="🛠️")

    content_slide("Database Design", [
        "## users table",
        "id (UUID PK), email (UNIQUE), password, role, created_at",
        "role is either 'student' or 'admin'",
        "---",
        "## feedback table",
        "id (UUID PK), student_id (FK → users.id)",
        "student_name, roll_no, phone_number",
        "subject, teacher, rating (1–5), comments",
        "admin_response, created_at",
        "---",
        "## Relationship",
        "One user → many feedback rows  (1-to-many)",
        "Demo users inserted via schema.sql seed data",
    ], icon="🗄️")

    content_slide("Feedback Submission (Students)", [
        "Students access feedback.html after logging in",
        "---",
        "## Form Fields",
        "Student Name  (text, required)",
        "Roll Number   (number, max 5 digits)",
        "Phone Number  (used for password reset verification)",
        "Subject       (text, required)",
        "Faculty Name  (text, required)",
        "Rating        (interactive 1–5 star selector)",
        "Comments      (text area, optional)",
        "---",
        "Data is sent via POST /api/feedback to Flask backend",
        "Backend saves to Supabase feedback table",
    ], icon="📝")

    content_slide("Admin Dashboard — Overview Tab", [
        "## 6 Live Stat Cards",
        "Total Feedbacks  |  Average Rating  |  Low-Rating Alerts",
        "Top Faculty  |  Most Reviewed Subject  |  Satisfaction Rate %",
        "---",
        "## Faculty Leaderboard",
        "All faculty ranked by average rating",
        "Top 3 highlighted with gold/silver/bronze medals",
        "Shows review count alongside average",
        "---",
        "## Recent Activity Feed",
        "Last 6 feedback submissions in reverse chronological order",
        "Colour-coded dots: green (good) / amber (average) / red (poor)",
        "Time-ago format (e.g. '5m ago', '2h ago')",
    ], icon="📊")

    content_slide("Admin Dashboard — Other Tabs", [
        "## 📈 Visual Analytics",
        "Bar Chart: Rating distribution (1★ to 5★ counts)",
        "Pie Chart: Feedback breakdown by subject",
        "---",
        "## 📋 Detailed Logs",
        "Full table of all feedback with admin response box",
        "Admin can reply to individual feedback entries",
        "Rows highlighted in yellow when navigated from notifications",
        "---",
        "## 🔍 Search Feedback",
        "Search by Student Roll Number for instant results",
        "---",
        "## 🎨 Theme Settings",
        "Toggle between Light Mode and Dark Mode",
    ], icon="🖥️")

    content_slide("Alerts & Notifications", [
        "Dedicated sidebar tab with red badge showing unread count",
        "---",
        "## 🚨 Alerts Panel",
        "Lists all feedback with rating below 2",
        "Shows student name, teacher rated, subject & comment",
        "Instantly highlights problematic faculty performance",
        "---",
        "## 🔔 Notifications Panel",
        "Shows latest 10 feedback submissions (newest first)",
        "Colour-coded icons: green / amber / red by rating",
        "---",
        "## Click-to-Navigate",
        "Clicking any item → switches to Detailed Logs tab",
        "Auto-scrolls to that exact row with yellow flash animation",
    ], icon="🔔")

    content_slide("Password Reset — OTP Flow", [
        "## Step 1: Identity Verification",
        "Student enters Email + Phone Number",
        "Backend verifies email in users table",
        "Backend verifies phone in feedback table",
        "---",
        "## Step 2: OTP Entry",
        "6-digit OTP generated and stored in server memory",
        "OTP displayed on screen for demo purposes",
        "Production: send via SMS / Email",
        "---",
        "## Step 3: New Password",
        "Student enters & confirms new password",
        "Password updated in users table identified by email",
        "OTP cleared from memory after successful reset",
    ], icon="🔑")

    content_slide("API Endpoints", [
        "POST  /api/login               — Authenticate user",
        "POST  /api/register            — Create new account",
        "GET   /api/stats               — All feedback + stats (admin)",
        "POST  /api/feedback            — Submit new feedback",
        "POST  /api/admin-response      — Admin replies to feedback",
        "POST  /api/forgot-password     — Email OTP request",
        "POST  /api/reset-password      — Reset via email OTP",
        "POST  /api/forgot-password-phone — Phone + email verify",
        "POST  /api/reset-password-phone  — Reset via phone OTP",
    ], icon="🔗")

    two_col_slide("User Roles & Demo Credentials",
        [
            "STUDENT ROLE",
            "Register and log in",
            "Submit feedback for any faculty",
            "Reset password via OTP",
            "",
            "Demo Login:",
            "Email:    student@demo.com",
            "Password: student123",
        ],
        [
            "ADMIN ROLE",
            "Full dashboard access",
            "View all feedback & analytics",
            "Respond to student feedback",
            "View alerts & notifications",
            "",
            "Demo Login:",
            "Email:    admin@demo.com",
            "Password: admin123",
        ], icon="👤"
    )

    content_slide("How to Run the Project", [
        "## Step 1 — Set up Supabase",
        "Create project at supabase.com",
        "Run schema.sql in the Supabase SQL Editor",
        "Copy Project URL and anon API key",
        "---",
        "## Step 2 — Configure Environment",
        "Edit backend/.env with your SUPABASE_URL and SUPABASE_KEY",
        "---",
        "## Step 3 — Install & Run Backend",
        "cd backend  →  pip install -r requirements.txt",
        "python app.py   (runs on http://localhost:5000)",
        "---",
        "## Step 4 — Open Frontend",
        "Open frontend/index.html in a browser",
        "Or: python -m http.server 8080 then visit localhost:8080",
    ], icon="▶️")

    content_slide("Future Enhancements", [
        "🔴 HIGH — Hash passwords using bcrypt before storing",
        "🟡 MEDIUM — Send OTP via real SMS (Twilio) or Email (SMTP)",
        "🟡 MEDIUM — Export feedback reports to PDF / Excel",
        "🟡 MEDIUM — Make frontend fully mobile responsive",
        "🟢 LOW — Faculty self-login to view their own ratings",
        "🟢 LOW — Filter feedback by semester / academic year",
        "🟢 LOW — Multi-language support (English / regional)",
        "---",
        "## Security Improvements",
        "Replace localStorage auth with JWT + httpOnly cookies",
        "Add server-side role verification middleware",
        "Rate-limit API endpoints to prevent abuse",
    ], icon="🚀")

    thank_you_slide()

    out_path = r"c:\Users\keert\OneDrive\Desktop\presentation\PRESENTATION.pptx"
    prs.save(out_path)
    print(f"✅ PowerPoint saved: {out_path}")


# ─── HTML → PDF Generation ──────────────────────────────────────────────────

def generate_pdf_html():
    html = """<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<title>Faculty Feedback System — Documentation</title>
<style>
  @import url('https://fonts.googleapis.com/css2?family=Outfit:wght@400;600;700;800&display=swap');

  * { margin: 0; padding: 0; box-sizing: border-box; }

  body {
    font-family: 'Outfit', sans-serif;
    background: #fff;
    color: #1e293b;
    font-size: 13px;
    line-height: 1.6;
  }

  /* ── Cover Page ─────────── */
  .cover {
    page-break-after: always;
    background: #1e293b;
    color: white;
    height: 100vh;
    display: flex;
    flex-direction: column;
    align-items: center;
    justify-content: center;
    text-align: center;
    padding: 60px;
    position: relative;
  }
  .cover-logo {
    width: 72px; height: 72px; margin-bottom: 28px;
  }
  .cover h1 {
    font-size: 42px; font-weight: 800; color: #fff; margin-bottom: 14px;
    background: linear-gradient(to right, #818cf8, #6366f1);
    -webkit-background-clip: text; -webkit-text-fill-color: transparent;
  }
  .cover .subtitle {
    font-size: 16px; color: #94a3b8; margin-bottom: 10px;
  }
  .cover .college {
    font-size: 14px; color: #6366f1; font-weight: 600; letter-spacing: 0.1em;
    text-transform: uppercase; margin-bottom: 6px;
  }
  .cover .date { font-size: 12px; color: #475569; }
  .cover-bar {
    position: absolute; bottom: 0; left: 0; right: 0;
    height: 6px; background: linear-gradient(to right, #6366f1, #4f46e5);
  }

  /* ── Pages ──────────────── */
  .page {
    page-break-after: always;
    padding: 40px 50px;
    min-height: 100vh;
  }
  .page:last-child { page-break-after: avoid; }

  /* ── Section Heading ─────── */
  .section-heading {
    font-size: 22px; font-weight: 800; color: #1e293b;
    border-left: 5px solid #6366f1;
    padding-left: 14px;
    margin-bottom: 24px;
  }

  /* ── Sub heading ─────────── */
  h3 {
    font-size: 13px; font-weight: 700; color: #6366f1;
    text-transform: uppercase; letter-spacing: 0.06em;
    margin: 18px 0 8px 0;
  }

  /* ── Table ───────────────── */
  table { width: 100%; border-collapse: collapse; margin: 16px 0; font-size: 12px; }
  th { background: #6366f1; color: white; padding: 8px 12px; text-align: left; font-weight: 600; }
  td { padding: 8px 12px; border-bottom: 1px solid #e2e8f0; }
  tr:nth-child(even) td { background: #f8fafc; }

  /* ── Code block ──────────── */
  pre {
    background: #1e293b; color: #e2e8f0;
    padding: 16px 20px; border-radius: 10px;
    font-size: 11px; overflow-x: auto;
    margin: 12px 0; line-height: 1.7;
    font-family: 'Courier New', monospace;
  }

  /* ── Stat card grid ──────── */
  .card-grid { display: grid; grid-template-columns: repeat(3,1fr); gap: 14px; margin: 16px 0; }
  .card {
    background: #f8fafc; border: 1px solid #e2e8f0;
    border-radius: 12px; padding: 16px;
    border-top: 4px solid #6366f1;
  }
  .card .icon { font-size: 22px; margin-bottom: 6px; }
  .card .label { font-size: 10px; color: #64748b; text-transform: uppercase; letter-spacing: .06em; font-weight: 600; }
  .card .value { font-size: 18px; font-weight: 800; color: #6366f1; }

  /* ── Step box ────────────── */
  .step { display: flex; gap: 14px; margin-bottom: 14px; align-items: flex-start; }
  .step-num {
    background: #6366f1; color: white;
    width: 28px; height: 28px; border-radius: 50%;
    display: flex; align-items: center; justify-content: center;
    font-weight: 700; font-size: 13px; flex-shrink: 0;
  }
  .step-body p { font-size: 12px; color: #64748b; margin-top: 2px; }

  /* ── Role cards ──────────── */
  .role-grid { display: grid; grid-template-columns: 1fr 1fr; gap: 16px; margin-top: 16px; }
  .role-card { background: #f8fafc; border-radius: 12px; padding: 20px; border: 1px solid #e2e8f0; }
  .role-card h4 { color: #6366f1; font-size: 14px; margin-bottom: 10px; }
  .role-card li { font-size: 12px; color: #475569; margin-left: 16px; margin-bottom: 4px; }
  .cred { background:#1e293b; color: #a5b4fc; padding: 8px 12px; border-radius: 8px; font-size: 11px; margin-top: 10px; font-family: monospace; }

  /* ── Footer ──────────────── */
  @media print {
    body { -webkit-print-color-adjust: exact; print-color-adjust: exact; }
    .cover { height: 100vh; }
    @page { margin: 0; size: A4; }
  }

  .footer {
    position: fixed; bottom: 0; left: 0; right: 0;
    background: white; border-top: 1px solid #e2e8f0;
    padding: 6px 50px; font-size: 10px; color: #94a3b8;
    display: flex; justify-content: space-between;
  }

  .tag {
    display: inline-block; background: #eef2ff; color: #6366f1;
    padding: 2px 10px; border-radius: 20px; font-size: 11px;
    font-weight: 600; margin: 2px;
  }
  .tag.green  { background: #ecfdf5; color: #10b981; }
  .tag.red    { background: #fef2f2; color: #ef4444; }
  .tag.amber  { background: #fffbeb; color: #f59e0b; }

  ul li { margin: 4px 0 4px 18px; color: #475569; }
</style>
</head>
<body>

<!-- COVER -->
<div class="cover">
  <svg class="cover-logo" viewBox="0 0 44 44" fill="none">
    <defs>
      <linearGradient id="lg" x1="0" y1="0" x2="44" y2="44" gradientUnits="userSpaceOnUse">
        <stop offset="0%" stop-color="#6366f1"/>
        <stop offset="100%" stop-color="#4f46e5"/>
      </linearGradient>
    </defs>
    <rect width="44" height="44" rx="12" fill="url(#lg)"/>
    <polygon points="22,10 36,17 22,24 8,17" fill="white" opacity="0.95"/>
    <path d="M12 18.5v6c0 2.5 4.5 5 10 5s10-2.5 10-5v-6" stroke="white" stroke-width="2" fill="none" stroke-linecap="round"/>
    <line x1="36" y1="17" x2="36" y2="24" stroke="white" stroke-width="2" stroke-linecap="round" opacity="0.8"/>
    <circle cx="36" cy="25.5" r="1.8" fill="white" opacity="0.8"/>
    <circle cx="33" cy="33" r="7" fill="#fbbf24"/>
    <path d="M33 28.5l1.2 3.6h3.8l-3 2.2 1.1 3.4-3.1-2.2-3.1 2.2 1.1-3.4-3-2.2h3.8z" fill="white"/>
  </svg>
  <h1>Faculty Feedback System</h1>
  <p class="subtitle">A digital platform for collecting &amp; analysing faculty feedback</p>
  <p class="college">Your Institution</p>
  <p class="date">Project Documentation &nbsp;|&nbsp; February 2026</p>
  <div class="cover-bar"></div>
</div>

<!-- PAGE 1 — Overview -->
<div class="page">
  <div class="section-heading">📌 Project Overview</div>
  <p>The <strong>Faculty Feedback System</strong> is a full-stack web application developed for your institution that allows students to submit feedback about faculty members. Administrators can view, analyse, and respond to feedback through a rich admin dashboard.</p>

  <h3>Key Objectives</h3>
  <ul>
    <li>Enable students to rate faculty on a <strong>1–5 star scale</strong> with comments</li>
    <li>Provide administrators a <strong>real-time analytics dashboard</strong></li>
    <li>Instant <strong>low-rating alerts</strong> to detect problematic patterns early</li>
    <li>Secure <strong>role-based access</strong> for students and admins</li>
    <li>OTP-based <strong>password reset</strong> via phone + email verification</li>
  </ul>

  <h3>Technology Stack</h3>
  <table>
    <tr><th>Layer</th><th>Technology</th></tr>
    <tr><td>Frontend</td><td>HTML5, CSS3, Vanilla JavaScript, Chart.js</td></tr>
    <tr><td>Backend</td><td>Python 3, Flask, Flask-CORS, python-dotenv</td></tr>
    <tr><td>Database</td><td>Supabase (PostgreSQL) — cloud hosted</td></tr>
    <tr><td>Fonts / UI</td><td>Google Fonts (Outfit), inline SVG icons</td></tr>
  </table>
</div>

<!-- PAGE 2 — Structure + DB -->
<div class="page">
  <div class="section-heading">🗄️ Project Structure &amp; Database</div>

  <h3>Folder Structure</h3>
  <pre>presentation/
├── backend/
│   ├── app.py            # All Flask API routes
│   ├── database.py       # Supabase client init
│   ├── requirements.txt  # Python dependencies
│   └── .env              # Supabase credentials
│
├── frontend/
│   ├── index.html        # Login page
│   ├── register.html     # Registration page
│   ├── feedback.html     # Feedback form (students)
│   ├── dashboard.html    # Admin dashboard
│   ├── forgot.html       # Password reset (OTP)
│   ├── css/style.css     # Global stylesheet
│   └── js/auth.js        # Auth logic
│
└── schema.sql            # Database setup script</pre>

  <h3>Database Tables</h3>
  <table>
    <tr><th>Column</th><th>Type</th><th>Description</th></tr>
    <tr><td colspan="3" style="background:#eef2ff;color:#4f46e5;font-weight:700;">users</td></tr>
    <tr><td>id</td><td>UUID PK</td><td>Auto-generated unique user ID</td></tr>
    <tr><td>email</td><td>TEXT UNIQUE</td><td>User's email address</td></tr>
    <tr><td>password</td><td>TEXT</td><td>User's password</td></tr>
    <tr><td>role</td><td>TEXT</td><td>'student' or 'admin'</td></tr>
    <tr><td colspan="3" style="background:#eef2ff;color:#4f46e5;font-weight:700;">feedback</td></tr>
    <tr><td>id</td><td>UUID PK</td><td>Auto-generated feedback ID</td></tr>
    <tr><td>student_id</td><td>UUID FK</td><td>References users.id</td></tr>
    <tr><td>student_name, roll_no</td><td>TEXT</td><td>Student identity</td></tr>
    <tr><td>phone_number</td><td>TEXT</td><td>Used for password reset verification</td></tr>
    <tr><td>subject, teacher</td><td>TEXT</td><td>Feedback target</td></tr>
    <tr><td>rating</td><td>INTEGER 1–5</td><td>Star rating</td></tr>
    <tr><td>comments, admin_response</td><td>TEXT</td><td>Written feedback &amp; reply</td></tr>
  </table>
</div>

<!-- PAGE 3 — Dashboard Features -->
<div class="page">
  <div class="section-heading">📊 Admin Dashboard Features</div>

  <div class="card-grid">
    <div class="card"><div class="icon">📝</div><div class="label">Total Feedbacks</div><div class="value">Live count</div></div>
    <div class="card"><div class="icon">⭐</div><div class="label">Average Rating</div><div class="value">Out of 5.0</div></div>
    <div class="card"><div class="icon">🚨</div><div class="label">Low-Rating Alerts</div><div class="value">Rating &lt; 2</div></div>
    <div class="card"><div class="icon">🏆</div><div class="label">Top Faculty</div><div class="value">Highest avg</div></div>
    <div class="card"><div class="icon">📚</div><div class="label">Top Subject</div><div class="value">Most reviewed</div></div>
    <div class="card"><div class="icon">😊</div><div class="label">Satisfaction Rate</div><div class="value">Rating ≥ 4 %</div></div>
  </div>

  <h3>Dashboard Tabs</h3>
  <table>
    <tr><th>Tab</th><th>What it shows</th></tr>
    <tr><td>📊 Overview</td><td>6 stat cards, faculty leaderboard 🥇🥈🥉, recent activity feed</td></tr>
    <tr><td>📈 Analytics</td><td>Bar chart (rating distribution) + Pie chart (by subject)</td></tr>
    <tr><td>📋 Detailed Logs</td><td>Full feedback table with admin response functionality</td></tr>
    <tr><td>🔍 Search</td><td>Search feedback by Student Roll Number</td></tr>
    <tr><td>🔔 Alerts &amp; Notifications</td><td>Alerts (rating &lt; 2) + latest 10 notifications, click-to-navigate</td></tr>
    <tr><td>🎨 Theme Settings</td><td>Light / Dark mode toggle</td></tr>
  </table>

  <h3>Click-to-Navigate Feature</h3>
  <p>Clicking any notification item automatically switches to the <strong>Detailed Logs</strong> tab, smoothly scrolls to the exact feedback row, and highlights it with a <strong>yellow flash animation</strong> for easy identification.</p>
</div>

<!-- PAGE 4 — API + Auth -->
<div class="page">
  <div class="section-heading">🔗 API Endpoints &amp; Authentication</div>

  <h3>REST API Endpoints  (http://localhost:5000)</h3>
  <table>
    <tr><th>Method</th><th>Endpoint</th><th>Description</th></tr>
    <tr><td><span class="tag">POST</span></td><td>/api/login</td><td>Authenticate user, return role</td></tr>
    <tr><td><span class="tag">POST</span></td><td>/api/register</td><td>Create new student/admin account</td></tr>
    <tr><td><span class="tag green">GET</span></td><td>/api/stats</td><td>All feedback + aggregate stats (admin)</td></tr>
    <tr><td><span class="tag">POST</span></td><td>/api/feedback</td><td>Submit new feedback (students)</td></tr>
    <tr><td><span class="tag">POST</span></td><td>/api/admin-response</td><td>Admin submits reply to feedback</td></tr>
    <tr><td><span class="tag amber">POST</span></td><td>/api/forgot-password-phone</td><td>Verify email + phone, issue OTP</td></tr>
    <tr><td><span class="tag amber">POST</span></td><td>/api/reset-password-phone</td><td>Reset password using phone OTP</td></tr>
  </table>

  <h3>Password Reset — 3-Step OTP Flow</h3>
  <div class="step">
    <div class="step-num">1</div>
    <div class="step-body"><strong>Identity Verification</strong><p>Student enters Email + Phone. Backend verifies email in users table and phone in feedback table.</p></div>
  </div>
  <div class="step">
    <div class="step-num">2</div>
    <div class="step-body"><strong>OTP Entry</strong><p>6-digit OTP generated server-side and stored in memory. Displayed on screen for demo purposes.</p></div>
  </div>
  <div class="step">
    <div class="step-num">3</div>
    <div class="step-body"><strong>New Password</strong><p>Student sets new password. Updated in users table by email. OTP cleared from memory.</p></div>
  </div>

  <h3>User Roles &amp; Demo Credentials</h3>
  <div class="role-grid">
    <div class="role-card">
      <h4>🎓 Student</h4>
      <ul>
        <li>Register &amp; log in</li>
        <li>Submit feedback for faculty</li>
        <li>Reset password via OTP</li>
      </ul>
      <div class="cred">Email: student@demo.com<br>Pass:  student123</div>
    </div>
    <div class="role-card">
      <h4>🛡️ Admin</h4>
      <ul>
        <li>Full dashboard access</li>
        <li>View all feedback &amp; analytics</li>
        <li>Reply to student feedback</li>
        <li>View alerts &amp; notifications</li>
      </ul>
      <div class="cred">Email: admin@demo.com<br>Pass:  admin123</div>
    </div>
  </div>
</div>

<!-- PAGE 5 — Setup + Future -->
<div class="page">
  <div class="section-heading">▶️ How to Run &amp; Future Plans</div>

  <h3>Setup Steps</h3>
  <div class="step">
    <div class="step-num">1</div>
    <div class="step-body"><strong>Supabase Setup</strong><p>Create a project at supabase.com → Run schema.sql in SQL Editor → Copy URL &amp; API key.</p></div>
  </div>
  <div class="step">
    <div class="step-num">2</div>
    <div class="step-body"><strong>Configure .env</strong><p>Edit <code>backend/.env</code> with your SUPABASE_URL and SUPABASE_KEY values.</p></div>
  </div>
  <div class="step">
    <div class="step-num">3</div>
    <div class="step-body"><strong>Install &amp; Run Backend</strong>
    <pre style="margin:6px 0;">cd backend
pip install -r requirements.txt
python app.py</pre></div>
  </div>
  <div class="step">
    <div class="step-num">4</div>
    <div class="step-body"><strong>Open Frontend</strong><p>Open <code>frontend/index.html</code> in a browser — or serve via <code>python -m http.server 8080</code>.</p></div>
  </div>

  <h3>Security Considerations</h3>
  <table>
    <tr><th>Area</th><th>Current (Demo)</th><th>Production Recommendation</th></tr>
    <tr><td>Passwords</td><td>Plain text</td><td>Hash with bcrypt</td></tr>
    <tr><td>OTP Delivery</td><td>Displayed on screen</td><td>Send via SMS (Twilio) or Email</td></tr>
    <tr><td>Auth Tokens</td><td>LocalStorage</td><td>JWT + httpOnly cookies</td></tr>
    <tr><td>CORS</td><td>All origins allowed</td><td>Restrict to frontend domain</td></tr>
  </table>

  <h3>Future Enhancements</h3>
  <p>
    <span class="tag red">High</span> Password hashing (bcrypt) &nbsp;
    <span class="tag amber">Medium</span> SMS/Email OTP delivery &nbsp;
    <span class="tag amber">Medium</span> Export to PDF/Excel &nbsp;
    <span class="tag amber">Medium</span> Mobile responsive design &nbsp;
    <span class="tag green">Low</span> Faculty self-login portal &nbsp;
    <span class="tag green">Low</span> Semester/year filters &nbsp;
    <span class="tag green">Low</span> Multi-language support
  </p>

  <br><br>
  <div style="text-align:center;padding:30px;background:#1e293b;border-radius:16px;color:white;margin-top:30px;">
    <div style="font-size:26px;font-weight:800;background:linear-gradient(to right,#818cf8,#6366f1);-webkit-background-clip:text;-webkit-text-fill-color:transparent;">Faculty Feedback System</div>
    <div style="color:#6366f1;font-size:12px;letter-spacing:.1em;text-transform:uppercase;margin-top:6px;">Your Institution &nbsp;|&nbsp; Admin Dashboard</div>
    <div style="color:#475569;font-size:11px;margin-top:8px;">For questions: admin@demo.com</div>
  </div>
</div>

</body>
</html>
"""
    out_path = r"c:\Users\keert\OneDrive\Desktop\presentation\DOCUMENTATION_PDF.html"
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(html)
    print(f"✅ PDF-ready HTML saved: {out_path}")
    print("   → Open it in Chrome/Edge → Ctrl+P → 'Save as PDF'")


# ─── Main ────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    print("\n🎓 Faculty Feedback System — Doc Generator\n")
    generate_pdf_html()
    generate_pptx()
    print("\n🎉 Done! Check your presentation folder.")
